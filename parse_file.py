#!/usr/bin/env python3
"""File parser for GP intake - extracts text from various file formats.
Usage: echo 'base64_content' | python3 parse_file.py <filename>
Outputs extracted text to stdout."""

import sys, os, tempfile, base64, json

# Read base64 from stdin
b64 = sys.stdin.read().strip()
if not b64:
    print("NO_INPUT")
    sys.exit(0)

filename = sys.argv[1] if len(sys.argv) > 1 else 'unknown.txt'
ext = os.path.splitext(filename)[1].lower()

# Decode and write to temp file
raw = base64.b64decode(b64)
tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
tmp.write(raw)
tmp_path = tmp.name
tmp.close()

text = ''

try:
    if ext == '.txt':
        with open(tmp_path, 'r', errors='replace') as f:
            text = f.read()

    elif ext == '.pdf':
        import subprocess
        r = subprocess.run(['pdftotext', tmp_path, '-'], stdout=subprocess.PIPE, stderr=subprocess.PIPE, universal_newlines=True, timeout=30)
        text = r.stdout or '[PDF 解析无内容]'

    elif ext in ('.docx', '.doc'):
        try:
            from docx import Document
            doc = Document(tmp_path)
            text = '\n'.join([p.text for p in doc.paragraphs])
            if not text.strip():
                text = '[Word 文档无段落文本]'
        except Exception as e:
            text = f'[Word 解析失败: {e}]'

    elif ext in ('.pptx', '.ppt'):
        try:
            from pptx import Presentation
            prs = Presentation(tmp_path)
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        texts.append(shape.text)
                if texts:
                    slides.append(f'--- 第{i}页 ---\n' + '\n'.join(texts))
            text = '\n\n'.join(slides) if slides else '[PPT 无文本内容]'
        except Exception as e:
            text = f'[PPT 解析失败: {e}]'

    elif ext in ('.xlsx', '.xls'):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(tmp_path, read_only=True, data_only=True)
            parts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    vals = [str(c) if c is not None else '' for c in row]
                    rows.append('\t'.join(vals))
                parts.append(f'=== 工作表: {sheet} ===\n' + '\n'.join(rows))
            text = '\n\n'.join(parts)
            if not text.strip():
                text = '[Excel 无内容]'
        except Exception as e:
            text = f'[Excel 解析失败: {e}]'

    else:
        text = f'[不支持的文件格式: {ext}]'

finally:
    os.unlink(tmp_path)

# Truncate if too long (DeepSeek context limit)
max_chars = 50000
if len(text) > max_chars:
    text = text[:max_chars] + f'\n\n... [内容过长，已截断至 {max_chars} 字符]'

print(text)
